# -*- coding: utf-8 -*-
"""Gera o PPT (6 paginas, uma por frente) replicando o layout do painel HTML.

- conteudo: frentes.json (extraido do painel) + overrides.json (memorando L2S v2)
- os textos ficam em tabelas nativas do PowerPoint, para facilitar a edicao
- menu de navegacao com hiperlinks internos entre as paginas

Tokens visuais e helpers de tabela vivem em deckstyle.py.
"""
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from deckstyle import (
    GREEN, GREEN_DEEP, BLACK, TEXT_DARK, TEXT_REG, TEXT_HELPER, NAV_IDLE,
    WHITE, GRAY_X_LIGHT, GRAY_LIGHT, GRAY_MEDIUM, STROKE, PLACEHOLDER,
    STATUS, FONT_NAME, CELL_MX,
    text_w, wrap, nlines, lh, split_bold,
    no_shadow, style_run, cell_box, cell_text, set_table_plain, set_row_heights,
)


def textbox(slide, x, y, w, h, text, size=9, color=TEXT_DARK, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, caps=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    style_run(p.add_run(), text.upper() if caps else text, size, color, bold)
    return tb


def rect(slide, x, y, w, h, fill=None, line=None, lw=0.75,
         shape=MSO_SHAPE.RECTANGLE, adj=None):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    no_shadow(s)
    if adj is not None:
        try: s.adjustments[0] = adj
        except Exception: pass
    tf = s.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return s


def slide_link(shape, to_slide):
    """hiperlink interno na caixa: clicar leva direto para a pagina da frente.

    O link fica no shape (nao no run) para o texto manter a cor e nao vir
    sublinhado como hiperlink.
    """
    shape.click_action.target_slide = to_slide


def new_table(slide, nrows, ncols, x, y, w, heights, widths):
    gf = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y),
                                Inches(w), Inches(sum(heights)))
    tbl = gf.table
    set_table_plain(tbl)
    set_row_heights(tbl, heights)
    for i, cw in enumerate(widths):
        tbl.columns[i].width = Inches(cw)
    return tbl


# ---------- geometria da pagina ----------
SW, SH = 13.3333, 7.5
ML = 0.42
CW = SW - 2 * ML
TOP = 0.30
BOTTOM = SH - 0.40

# corpo maior; so reduz se a pagina exigir
BODY_STEPS = (10.0, 9.5, 9.0, 8.5, 8.0, 7.5)
PT_TITLE, PT_OBJ = 15.0, 10.5
PT_LABEL, PT_VALUE = 8.0, 10.0
PT_NAV, PT_SECTION, PT_FOOT = 9.5, 11.5, 8.0
PT_HEAD = 9.0

COLS = [
    ('#',                  0.62, PP_ALIGN.LEFT),
    ('Milestone',          3.30, PP_ALIGN.LEFT),
    ('Responsável',        1.42, PP_ALIGN.LEFT),
    ('Prazo',              0.95, PP_ALIGN.LEFT),
    ('Status',             1.15, PP_ALIGN.CENTER),
    ('Progresso recente',  1.68, PP_ALIGN.LEFT),
    ('Próximos passos',    1.68, PP_ALIGN.LEFT),
    ('Pontos a escalar',   1.68, PP_ALIGN.LEFT),
]
_scale = CW / sum(c[1] for c in COLS)
COLS = [(n, w * _scale, a) for n, w, a in COLS]
CELL_MX = 0.07
BAR = 0.075

# ---------- conteudo ----------
frentes = json.load(open('frentes.json'))
ovr = json.load(open('overrides.json'))
for f in frentes:
    fo = ovr['frentes'].get(str(f['id']), {})
    f['obj_bold'] = fo.get('obj_bold', [])
    f['nota'] = fo.get('nota', '')
    for m in f['entregaveis']:
        mo = ovr['milestones'].get(m['id'], {})
        m['desc'] = mo.get('desc', m['desc'])
        m['prazo'] = mo.get('prazo', m['prazo'])
        m['bold'] = mo.get('bold', [])

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
blank = prs.slide_layouts[6]
slides = [prs.slides.add_slide(blank) for _ in frentes]   # todas antes, p/ os links

for slide, f in zip(slides, frentes):
    rect(slide, 0, 0, SW, SH, fill=WHITE)
    y = TOP

    # ======== CABECALHO DA FRENTE (tabela: barra | titulo+objetivo | lider | time) ========
    w_lid, w_time = 1.90, 2.90
    w_main = CW - BAR - w_lid - w_time
    BADGE, BADGE_GAP = 0.30, 0.14
    MAIN_ML = BADGE + BADGE_GAP + 0.20

    txt_w = w_main - MAIN_ML - CELL_MX
    tit_n = nlines(f['nome'], PT_TITLE, txt_w, True)
    obj_n = nlines(f['objetivo'], PT_OBJ, txt_w)
    h_main = tit_n * lh(PT_TITLE) + 0.06 + obj_n * lh(PT_OBJ)
    side_n = max(nlines(f['lideres'], PT_VALUE, w_lid - 2 * CELL_MX),
                 nlines(f['time'], PT_VALUE, w_time - 2 * CELL_MX))
    h_side = lh(PT_LABEL) + 0.05 + side_n * lh(PT_VALUE)
    hdr_h = max(h_main, h_side) + 0.24

    hdr = new_table(slide, 1, 4, ML, y, CW, [hdr_h], [BAR, w_main, w_lid, w_time])
    cell_box(hdr.cell(0, 0), GREEN)
    for c in (1, 2, 3):
        cell_box(hdr.cell(0, c), GRAY_X_LIGHT)

    c = hdr.cell(0, 1)
    tf = cell_text(c, [(f['nome'], True)], PT_TITLE, BLACK, mt=0.12, mb=0.12)
    c.margin_left = Inches(MAIN_ML)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    for seg, b in split_bold(f['objetivo'], f['obj_bold']):
        style_run(p2.add_run(), seg, PT_OBJ, TEXT_REG, b)

    for col, lab, val in ((2, 'Líder(es)', f['lideres']),
                          (3, 'Time de trabalho', f['time'])):
        c = hdr.cell(0, col)
        tf = cell_text(c, [(lab.upper(), True)], PT_LABEL, TEXT_HELPER, mt=0.12, mb=0.12)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(3)
        style_run(p2.add_run(), val, PT_VALUE, TEXT_DARK)

    y += hdr_h + 0.13

    # ======== MENU DE FRENTES (hiperlinks internos) ========
    nav_pt = PT_NAV
    labels = ['%d. %s' % (g['id'], g['nome']) for g in frentes]
    while nav_pt > 6.5:
        ws = [text_w(t, nav_pt, True) for t in labels]
        gap = (CW - sum(ws)) / (len(labels) - 1)
        if gap >= 0.16:
            break
        nav_pt -= 0.5
    nav_h = lh(nav_pt) + 0.06
    xx = ML
    for g, tgt, wgt in zip(frentes, slides, ws):
        active = g['id'] == f['id']
        tb = textbox(slide, xx, y, wgt + 0.04, nav_h, '', size=nav_pt)
        run = tb.text_frame.paragraphs[0].add_run()
        style_run(run, '%d. %s' % (g['id'], g['nome']), nav_pt,
                  GREEN_DEEP if active else NAV_IDLE, bold=True)
        slide_link(tb, tgt)
        xx += wgt + gap
    y += nav_h + 0.14

    # ======== TABELA DE MILESTONES ========
    textbox(slide, ML, y, 3.0, 0.20, 'Milestones', size=PT_SECTION, color=BLACK, bold=True)
    if f['nota']:
        textbox(slide, ML + 3.0, y + 0.03, CW - 3.0, 0.18, f['nota'],
                size=PT_FOOT, color=TEXT_HELPER, align=PP_ALIGN.RIGHT)
    y += lh(PT_SECTION) + 0.10

    rows = f['entregaveis']
    avail = BOTTOM - y

    for body_pt in BODY_STEPS:
        id_pt = body_pt - 0.5
        head_h = lh(PT_HEAD) + 0.14
        heights = []
        for r in rows:
            n = max(nlines(r['id'],    id_pt,   COLS[0][1] - 2 * CELL_MX, True),
                    nlines(r['desc'],  body_pt, COLS[1][1] - 2 * CELL_MX),
                    nlines(r['owner'], body_pt, COLS[2][1] - 2 * CELL_MX),
                    nlines(r['prazo'], body_pt, COLS[3][1] - 2 * CELL_MX))
            heights.append(max(0.36, n * lh(body_pt) + 0.14))
        if head_h + sum(heights) <= avail:
            break

    # distribui o espaco que sobra entre as linhas
    caps = [max(h, 1.10) for h in heights]
    for _ in range(60):
        slack = avail - (head_h + sum(heights))
        movable = [i for i, h in enumerate(heights) if h < caps[i] - 1e-4]
        if slack <= 0.01 or not movable:
            break
        step = slack / len(movable)
        for i in movable:
            heights[i] = min(caps[i], heights[i] + step)

    tbl = new_table(slide, len(rows) + 1, len(COLS), ML, y, CW,
                    [head_h] + heights, [c[1] for c in COLS])

    for i, (name, _w, align) in enumerate(COLS):
        c = tbl.cell(0, i)
        cell_box(c, None, (None, None, None, BLACK), lw=1.0)
        cell_text(c, [(name, True)], PT_HEAD, TEXT_HELPER, align=align,
                  anchor=MSO_ANCHOR.BOTTOM, mt=0.02, mb=0.05)

    for ri, r in enumerate(rows, start=1):
        sep = GRAY_MEDIUM if ri == len(rows) else GRAY_LIGHT
        base = (None, None, None, sep)

        c = tbl.cell(ri, 0)
        cell_box(c, None, base, lw=0.5)
        cell_text(c, [(r['id'], True)], id_pt, TEXT_HELPER)

        for i, key, col in ((1, 'desc', TEXT_DARK), (2, 'owner', TEXT_DARK),
                            (3, 'prazo', TEXT_REG)):
            c = tbl.cell(ri, i)
            cell_box(c, None, base, lw=0.5)
            parts = split_bold(r[key], r['bold']) if key == 'desc' else [(r[key], False)]
            cell_text(c, parts, body_pt, col)

        lbl, bg, fg = STATUS[r['status']]
        c = tbl.cell(ri, 4)
        cell_box(c, bg, (WHITE, WHITE, WHITE, WHITE), lw=3.0)
        cell_text(c, [(lbl, True)], min(body_pt, 8.5), fg, align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE, mx=0.03)

        for i in (5, 6, 7):
            c = tbl.cell(ri, i)
            cell_box(c, WHITE, (GRAY_MEDIUM,) * 4, lw=0.5)
            cell_text(c, [('TBD', False)], body_pt, PLACEHOLDER)

    # badge numerada por cima da barra verde (desenhada depois das tabelas)
    bx = ML + BAR + 0.20
    by = TOP + 0.12 + (lh(PT_TITLE) - BADGE) / 2.0
    rect(slide, bx, by, BADGE, BADGE, fill=GREEN, shape=MSO_SHAPE.OVAL)
    textbox(slide, bx, by + 0.055, BADGE, 0.22, str(f['id']),
            size=12, color=BLACK, bold=True, align=PP_ALIGN.CENTER)

    # ---- rodape ----
    ln = slide.shapes.add_connector(1, Inches(ML), Inches(SH - 0.32),
                                    Inches(ML + CW), Inches(SH - 0.32))
    ln.line.color.rgb = STROKE
    ln.line.width = Pt(0.5)
    no_shadow(ln)
    textbox(slide, ML, SH - 0.27, CW * 0.8, 0.16,
            'Fonte: Memorando estratégico Lead to Sales — Aceleração Seminovos',
            size=PT_FOOT, color=TEXT_HELPER)
    textbox(slide, ML + CW - 1.0, SH - 0.27, 1.0, 0.16, str(f['id']),
            size=PT_FOOT, color=TEXT_HELPER, align=PP_ALIGN.RIGHT)

prs.core_properties.title = 'Lead-to-Sales — Frentes e Milestones'
prs.save('frentes_milestones.pptx')
print('ok')
