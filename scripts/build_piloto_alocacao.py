# -*- coding: utf-8 -*-
"""Pagina 'Piloto de alocacao de leads em loja proxima' (pre vs pos piloto).

Replica o layout da pagina original (cabecalhos chevron, tres paineis por grupo,
graficos marimekko com linha de media e delta em p.p.), com os wordings
migrados de 'loja escolhida / alocada' para 'loja proxima' e os dados
substituidos pelos paineis novos.
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- tokens
GREEN = RGBColor(0x10, 0x4C, 0x3E)   # Forest 1  - alocado em loja proxima
GOLD  = RGBColor(0xC6, 0xAA, 0x3D)   # Sunset 2  - NAO alocado em loja proxima
RED   = RGBColor(0xCC, 0x00, 0x00)   # Bain red
RUBY  = RGBColor(0x99, 0x00, 0x00)   # negativo
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY1 = RGBColor(0x33, 0x33, 0x33)
GREY2 = RGBColor(0x5C, 0x5C, 0x5C)
GREY3 = RGBColor(0x85, 0x85, 0x85)
GREY4 = RGBColor(0xB4, 0xB4, 0xB4)
GOLDD = RGBColor(0xAB, 0x89, 0x33)   # Sunset 1 - texto sobre fundo claro
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Arial"
LANG = "pt-BR"

# ---------------------------------------------------------------- grid (in)
SW, SH = 13.333, 7.5
TITLE_X, TITLE_Y, TITLE_W, TITLE_H = 0.30, 0.16, 10.9, 0.70
RULE_Y = 0.92
CHEV_Y, CHEV_H = 1.00, 0.44
LBL_X, LBL_W = 0.26, 0.52
PANEL_X, PANEL_W = 0.86, 12.20
PANEL_TOP, PANEL_H, PANEL_GAP = 1.54, 1.74, 0.10
NOTE_Y = 7.06

CHEV_L = (PANEL_X, 5.62)              # 0.86 -> 6.48
CHEV_R = (7.02, PANEL_W + PANEL_X - 7.02)  # 7.02 -> 13.06

# per-chart geometry: (gutter_x, bars_x, bars_w, media_x, media_w)
PRE = (0.94, 1.52, 4.38, 5.92, 0.56)
POS = (7.08, 7.64, 4.52, 12.18, 0.62)
DELTA_X, DELTA_W = 6.55, 0.50         # anotacao +/-~X% entre pre e pos

BAR_MAX_H = 0.80
BASE_DY = 1.24                        # base das barras, a partir do topo do painel

# ---------------------------------------------------------------- dados
CAT_A = "Alocado em loja próxima"
CAT_B = "NÃO alocado em loja próxima"

ROWS = [
    dict(
        label="Lojas do piloto¹", label_color=RED, border=RED, dashed_red=True,
        change="-~58%", change_color=RUBY, axis=1.20,
        pre=dict(media=0.5792, media_lbl="Média 0,58%", delta="-0,53 p.p.", delta_neg=True,
                 a=dict(conv=0.5534, lbl="0,55%", pct="95,2%", leads="7.228", w=95.16),
                 b=dict(conv=1.0870, lbl="1,09%", pct="4,8%",  leads="368",   w=4.84)),
        pos=dict(media=0.2460, media_lbl="Média 0,25%", delta=None, delta_neg=True,
                 a=dict(conv=0.1849, lbl="0,18%", pct="99,8%", leads="1.623", w=99.82),
                 b=dict(conv=33.33,  lbl="33,33%", pct="0,2%", leads="3", w=0.18, offscale=True)),
    ),
    dict(
        label="Grupo controle²", label_color=GREEN, border=GREY4, dashed_red=False,
        change="+~9%", change_color=GREEN, axis=0.80,
        pre=dict(media=0.5117, media_lbl="Média 0,51%", delta="-0,17 p.p.", delta_neg=True,
                 a=dict(conv=0.50315, lbl="0,50%", pct="95,0%", leads="70.555", w=95.02),
                 b=dict(conv=0.67568, lbl="0,68%", pct="5,0%",  leads="3.700",  w=4.98)),
        pos=dict(media=0.5587, media_lbl="Média 0,56%", delta="-0,08 p.p.", delta_neg=True,
                 a=dict(conv=0.55548, lbl="0,56%", pct="95,9%", leads="14.582", w=95.85),
                 b=dict(conv=0.63391, lbl="0,63%", pct="4,2%",  leads="631",    w=4.15)),
    ),
    dict(
        label="Restante do Brasil", label_color=GREEN, border=GREY4, dashed_red=False,
        change="+~23%", change_color=GREEN, axis=0.80,
        pre=dict(media=0.4580, media_lbl="Média 0,46%", delta="+0,05 p.p.", delta_neg=False,
                 a=dict(conv=0.46314, lbl="0,46%", pct="89,6%", leads="368.367", w=89.56),
                 b=dict(conv=0.41443, lbl="0,41%", pct="10,4%", leads="42.950",  w=10.44)),
        pos=dict(media=0.5641, media_lbl="Média 0,56%", delta="+0,06 p.p.", delta_neg=False,
                 a=dict(conv=0.56957, lbl="0,57%", pct="90,5%", leads="76.667", w=90.49),
                 b=dict(conv=0.50868, lbl="0,51%", pct="9,5%",  leads="8.060",   w=9.51)),
    ),
]

TITLE = ("Conversão média das lojas do piloto caiu ~58% pós piloto vs. alta de ~9% no grupo "
         "controle e ~23% no restante do Brasil, com base pós piloto ainda pequena")

CALLOUT = ("Conversão dos leads alocados em loja próxima caiu de 0,55% para 0,18% pós piloto "
           "(3 vendas em 1.623 leads)")

NOTE = ("Nota: (1) VCFTZ, VCGOI, VCBSI, VCREC, VCRAG; (2) lojas em até 15 km das lojas do piloto; "
        "(3) leads avaliáveis = leads com loja próxima identificável; médias são a razão das somas "
        "(total de vendas / total de leads); (4) pré piloto = 103 dias, pós piloto = 25 dias; vendas no "
        "pós piloto: 4 (lojas do piloto), 85 (grupo controle), 478 (restante do Brasil). "
        "Fonte: base de leads gerados, base de vendas faturadas")

# ---------------------------------------------------------------- helpers
def _lang(run):
    run.font._rPr.set('lang', LANG)

def txt(slide, x, y, w, h, parts, size=7, color=BLACK, bold=False,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=None, wrap=True,
        line=None):
    """parts: str or list of (text, {bold,underline,color,size}) tuples."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line:
        p.line_spacing = line
    if isinstance(parts, str):
        parts = [(parts, {})]
    for t, o in parts:
        if t == "\n":                      # nova linha = novo paragrafo
            p = tf.add_paragraph()
            p.alignment = align
            if line:
                p.line_spacing = line
            continue
        r = p.add_run()
        r.text = t
        f = r.font
        f.name = FONT
        f.size = Pt(o.get("size", size))
        f.bold = o.get("bold", bold)
        f.underline = o.get("underline", False)
        f.color.rgb = o.get("color", color)
        if spacing:
            f._rPr.set('spc', str(int(spacing * 100)))
        _lang(r)
    return tb

def rect(slide, x, y, w, h, fill=None, line_color=None, line_w=0.75,
         dash=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w)
        if dash:
            ln = s.line._get_or_add_ln()
            for e in ln.findall(qn('a:prstDash')):
                ln.remove(e)
            d = ln.makeelement(qn('a:prstDash'), {'val': dash})
            ln.append(d)
    s.text_frame.word_wrap = False
    return s

def hline(slide, x1, y, x2, color=GREY4, w=0.75, dash=None):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y),
                                   Inches(x2), Inches(y))
    c.line.color.rgb = color
    c.line.width = Pt(w)
    if dash:
        ln = c.line._get_or_add_ln()
        d = ln.makeelement(qn('a:prstDash'), {'val': dash})
        ln.append(d)
    return c

def vline(slide, x, y1, y2, color=GREY2, w=0.75, heads=True):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y1),
                                   Inches(x), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(w)
    if heads:
        ln = c.line._get_or_add_ln()
        for tag, attrs in (('a:headEnd', {'type': 'triangle', 'w': 'sm', 'len': 'sm'}),
                           ('a:tailEnd', {'type': 'triangle', 'w': 'sm', 'len': 'sm'})):
            ln.append(ln.makeelement(qn(tag), attrs))
    return c

def rot_label(slide, x, y, w, h, text, color):
    tb = txt(slide, x, y, w, h, text, size=9, color=color, bold=True,
             align=PP_ALIGN.CENTER, wrap=True, line=0.9)
    tb.rotation = 270
    return tb

# ---------------------------------------------------------------- chart
def draw_chart(slide, geom, T, axis, d):
    gut_x, bx, bw, mx, mw = geom
    base = T + BASE_DY

    def bar_h(v):
        return max(min(v / axis, 1.0) * BAR_MAX_H, 0.012)

    # larguras marimekko proporcionais ao % dos leads
    tot = d["a"]["w"] + d["b"]["w"]
    wa = max(bw * d["a"]["w"] / tot, 0.03)
    wb = max(bw * d["b"]["w"] / tot, 0.022)
    xa, xb = bx, bx + bw - wb

    # eixo
    hline(slide, bx, base, bx + bw, GREY1, 0.75)

    OFF_H = BAR_MAX_H * 0.60          # barra fora de escala nao vai ao topo
    bars = []
    for spec, x, w, col in ((d["a"], xa, wa, GREEN), (d["b"], xb, wb, GOLD)):
        off = spec.get("offscale", False)
        h = OFF_H if off else bar_h(spec["conv"])
        rect(slide, x, base - h, w, h, fill=col)
        bars.append((spec, x, w, h, off))

    # linha de media (vermelha tracejada) + rotulo a direita
    hm = bar_h(d["media"])
    hline(slide, bx, base - hm, bx + bw, RED, 0.75, dash='dash')
    txt(slide, mx, base - hm - 0.075, mw, 0.15, d["media_lbl"], size=6.5,
        color=RED, bold=True, align=PP_ALIGN.LEFT)

    # rotulos de valor sobre as barras.
    # Sempre acima do maior entre topo da barra e linha de media, para o rotulo
    # nunca cair sobre a tracejada vermelha (linhas 1 e 3 empatam barra e media).
    (sa, xa_, wa_, ha, _), (sb, xb_, wb_, hb, off_b) = bars
    txt(slide, xa_ + wa_ / 2 - 0.35, min(base - ha, base - hm) - 0.155, 0.70, 0.14,
        sa["lbl"], size=7, color=GREY1)
    if off_b:
        # seta vermelha acima da barra sinalizando valor fora da escala do eixo
        c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       Inches(xb_ + wb_ / 2), Inches(base - hb + 0.02),
                                       Inches(xb_ + wb_ / 2), Inches(base - hb - 0.14))
        c.line.color.rgb = RED
        c.line.width = Pt(1.25)
        ln = c.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'sm', 'len': 'sm'}))
        txt(slide, xb_ - 0.79, base - hb - 0.09, 0.76, 0.14, sb["lbl"],
            size=7, color=RED, bold=True, align=PP_ALIGN.RIGHT)
    else:
        # barra B e estreita: rotulo colado no seu limite direito, e afastado o
        # bastante para nao encostar no rotulo de delta quando a seta e curta
        gap = abs(hb - hm)
        lift = 0.20 if gap < 0.18 else 0.155
        txt(slide, xb_ + wb_ - 0.40, min(base - hb, base - hm) - lift, 0.40, 0.14,
            sb["lbl"], size=7, color=GREY1, align=PP_ALIGN.RIGHT)

    # delta em p.p.: seta vertical da media ao topo da barra B, com o rotulo numa
    # pilula branca a esquerda da barra (mesmo padrao dos paineis de origem)
    if d["delta"]:
        y1, y2 = base - hm, base - hb
        col = RUBY if d["delta_neg"] else GREY2
        vline(slide, xb_ - 0.015, min(y1, y2), max(y1, y2), col, 0.75)
        rect(slide, xb_ - 0.47, (y1 + y2) / 2 - 0.07, 0.44, 0.14, fill=WHITE)
        txt(slide, xb_ - 0.47, (y1 + y2) / 2 - 0.07, 0.44, 0.14, d["delta"], size=6.5,
            color=col, bold=True, align=PP_ALIGN.RIGHT)

    # rotulos de categoria + linhas de % dos leads e # leads
    y_cat, y_pct, y_led = base + 0.02, base + 0.185, base + 0.315
    txt(slide, gut_x, y_pct, 0.58, 0.135, "% dos leads", size=6, color=GREY2,
        align=PP_ALIGN.LEFT, wrap=False)
    txt(slide, gut_x, y_led, 0.58, 0.135, "# leads", size=6, color=GREY2,
        align=PP_ALIGN.LEFT, wrap=False)

    txt(slide, xa_ + wa_ / 2 - 0.95, y_cat, 1.90, 0.15, CAT_A, size=6, color=GREEN, bold=True)
    txt(slide, xa_ + wa_ / 2 - 0.50, y_pct, 1.00, 0.135, sa["pct"], size=6, color=GREY1)
    txt(slide, xa_ + wa_ / 2 - 0.50, y_led, 1.00, 0.135, sa["leads"], size=6, color=GREY1)

    x_bl = bx + bw - 1.95
    txt(slide, x_bl, y_cat, 1.95, 0.15, CAT_B, size=6, color=GOLDD,
        bold=True, align=PP_ALIGN.RIGHT)
    txt(slide, bx + bw - 1.00, y_pct, 1.00, 0.135, sb["pct"], size=6, color=GREY1,
        align=PP_ALIGN.RIGHT)
    txt(slide, bx + bw - 1.00, y_led, 1.00, 0.135, sb["leads"], size=6, color=GREY1,
        align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------- build
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# titulo + tag PRELIMINAR + regua vermelha
txt(slide, TITLE_X, TITLE_Y, TITLE_W, TITLE_H, TITLE, size=19, color=GREY1,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=1.08)
txt(slide, 11.28, TITLE_Y + 0.40, 1.78, 0.24, "/ P R E L I M I N A R", size=8.5,
    color=GREY3, align=PP_ALIGN.RIGHT, wrap=False)
hline(slide, TITLE_X, RULE_Y, PANEL_X + PANEL_W, RED, 2.25)

# cabecalhos chevron
for (cx, cw), fill, per in ((CHEV_L, BLACK, "PRÉ PILOTO (01/04 a 13/07)"),
                            (CHEV_R, GREY2, "PÓS PILOTO (13/07 a 06/08)")):
    rect(slide, cx, CHEV_Y, cw, CHEV_H, fill=fill, shape=MSO_SHAPE.PENTAGON)
    txt(slide, cx + 0.12, CHEV_Y + 0.03, cw - 0.55, CHEV_H - 0.06,
        [("% conversão – leads com ", {}), ("loja próxima", {"underline": True}),
         (" identificada", {}), ("\n", {}), (per, {"bold": True, "size": 10})],
        size=8.5, color=WHITE, line=1.0)

# paineis
for i, row in enumerate(ROWS):
    T = PANEL_TOP + i * (PANEL_H + PANEL_GAP)
    rect(slide, PANEL_X, T, PANEL_W, PANEL_H, fill=None,
         line_color=row["border"], line_w=1.0 if row["dashed_red"] else 0.75, dash='dash')
    rot_label(slide, LBL_X - (PANEL_H - LBL_W) / 2, T + (PANEL_H - LBL_W) / 2,
              PANEL_H, LBL_W, row["label"], row["label_color"])
    draw_chart(slide, PRE, T, row["axis"], row["pre"])
    draw_chart(slide, POS, T, row["axis"], row["pos"])
    txt(slide, DELTA_X, T + 0.30, DELTA_W, 0.18, row["change"], size=9,
        color=row["change_color"], bold=True)

# callout da linha 1
T1 = PANEL_TOP
rect(slide, 8.00, T1 + 0.24, 3.22, 0.44, fill=WHITE, line_color=RED, line_w=1.0)
txt(slide, 8.06, T1 + 0.25, 3.10, 0.42, CALLOUT, size=6.5, color=RED, line=1.05)

# nota / fonte
txt(slide, TITLE_X, NOTE_Y, PANEL_X + PANEL_W - TITLE_X, 0.34, NOTE, size=6,
    color=GREY3, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=1.15)

OUT = "/home/user/Gerais/output/piloto_alocacao_loja_proxima.pptx"
prs.save(OUT)
print("saved", OUT)
