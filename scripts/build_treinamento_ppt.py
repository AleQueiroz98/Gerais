#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Monta o PPT (17 paginas, 16:9) do material de treinamento do Simulador de
Financiamento do PDV, replicando pixel a pixel o HTML de origem.

O layout vem medido do Chromium (extract_treinamento_layout.js): cada bloco do
HTML virou forma nativa e cada texto virou caixa de texto do PowerPoint, com as
mesmas coordenadas, cores, corpos de fonte e entrelinhas do original -- tudo
editavel, sem imagem de fundo.

    node extract_treinamento_layout.js ../output/treinamento_simulador_pdv.html \
         layout.json telas/
    python3 build_treinamento_ppt.py layout.json telas/ ../output/xxx.pptx
"""
import json
import pathlib
import re
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from PIL import Image
from pptx.util import Emu, Pt

from deckstyle import cell_border, no_shadow, set_row_heights, set_table_plain

# ---------- unidades ----------
PXE = 9525                      # 1 px CSS (96 dpi) = 9525 EMU
SW, SH = 1280, 720              # palco do HTML = 16:9 exato
FONT = 'Segoe UI'               # mesma familia declarada no HTML
FONT_MEDIDA = 'Liberation Sans'  # fonte usada na medicao do layout (o
                                 # container nao tem Segoe UI instalada)
# (ascender + |descender|) / em de cada fonte: e a "area de conteudo" que o
# CSS usa para distribuir a meia-entrelinha e de onde o PowerPoint tira a
# primeira linha de base (topo + entrelinha - descendente)
CONTEUDO = {'Segoe UI': 1.3301, 'Liberation Sans': 1.1172}
LANG = 'pt-BR'
ALIGN = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
         'right': PP_ALIGN.RIGHT, 'justify': PP_ALIGN.JUSTIFY,
         'start': PP_ALIGN.LEFT, 'end': PP_ALIGN.RIGHT}

GRADIENTE = {                   # linear-gradient(135deg, ...) das paginas cheias
    'capa': (45, [(0.0, '00602F'), (0.55, '004A25'), (1.0, '00351B')]),
    'fim':  (45, [(0.0, '00602F'), (1.0, '004A25')]),
}


def emu(px):
    return Emu(int(round(px * PXE)))


def pt(px):
    return Pt(px * 0.75)        # 1 px CSS = 0.75 pt


def num(v, default=0.0):
    try:
        return float(str(v).replace('px', ''))
    except (TypeError, ValueError):
        return default


def topo_linha(y_texto, corpo, linha):
    """topo da caixa de linha do CSS a partir do topo do texto medido"""
    if not linha:
        return y_texto
    return y_texto - (linha - CONTEUDO[FONT_MEDIDA] * corpo) / 2


def topo_caixa(y_linha, corpo, linha):
    """onde ancorar a caixa de texto para a linha de base cair no mesmo lugar
    do HTML: o CSS centra a area de conteudo na entrelinha, o PowerPoint mede
    a primeira linha de base a partir do topo."""
    if not linha:
        return y_linha
    return y_linha + (CONTEUDO[FONT] * corpo - linha) / 2


def color(css):
    """'rgb(a)(...)' -> (RGBColor, alpha) | (None, 0) quando transparente"""
    if not css:
        return None, 0.0
    m = re.match(r'rgba?\(([^)]+)\)', css)
    if not m:
        return None, 0.0
    p = [x.strip() for x in m.group(1).replace('/', ',').split(',')]
    r, g, b = (int(round(float(x))) for x in p[:3])
    a = float(p[3]) if len(p) > 3 else 1.0
    return RGBColor(r, g, b), a


def hexcolor(css):
    c, _ = color(css)
    return c


# ---------- formas ----------
def _alpha(fill_el, a):
    """alfa do CSS no solidFill (a:srgbClr/a:alpha)"""
    if a >= 1:
        return
    clr = fill_el.find(qn('a:srgbClr'))
    el = clr.makeelement(qn('a:alpha'), {'val': str(int(round(a * 100000)))})
    clr.append(el)


def rect(slide, x, y, w, h, fill=None, alpha=1.0, line=None, lw=0.0,
         line_alpha=1.0, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, emu(x), emu(y), emu(max(w, 0.1)),
                               emu(max(h, 0.1)))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        _alpha(s.fill._xPr.find(qn('a:solidFill')), alpha)
    if line is None or lw <= 0:
        s.line.fill.background()
    else:
        s.line.fill.solid()
        s.line.fill.fore_color.rgb = line
        s.line.width = pt(lw)
        _alpha(s.line._get_or_add_ln().find(qn('a:solidFill')), line_alpha)
    no_shadow(s)
    tf = s.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    return s


def grad_rect(slide, x, y, w, h, ang, stops):
    """retangulo com gradiente linear (angulo ja convertido para DrawingML)"""
    s = rect(slide, x, y, w, h, fill=RGBColor(0, 0, 0))
    spPr = s._element.spPr
    for e in spPr.findall(qn('a:solidFill')):
        spPr.remove(e)
    g = spPr.makeelement(qn('a:gradFill'), {'flip': 'none', 'rotWithShape': '1'})
    lst = g.makeelement(qn('a:gsLst'), {})
    for pos, hx in stops:
        gs = lst.makeelement(qn('a:gs'), {'pos': str(int(round(pos * 100000)))})
        gs.append(gs.makeelement(qn('a:srgbClr'), {'val': hx}))
        lst.append(gs)
    g.append(lst)
    g.append(g.makeelement(qn('a:lin'), {'ang': str(int(ang * 60000)),
                                         'scaled': '0'}))
    ln = spPr.find(qn('a:ln'))
    (ln.addprevious(g) if ln is not None else spPr.append(g))
    return s


def triangle(slide, x, y, w, h, fill):
    """seta do ciclo: triangulo apontando para a direita (forma editavel)"""
    ff = slide.shapes.build_freeform(emu(x), emu(y), scale=1.0)
    ff.add_line_segments([(emu(x + w), emu(y + h / 2)), (emu(x), emu(y + h))],
                         close=True)
    s = ff.convert_to_shape()
    s.shadow.inherit = False
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    no_shadow(s)
    return s


def shadow(shape, dist_px, blur_px, alpha):
    """box-shadow 0 Ypx Zpx rgba(0,0,0,a) -> outerShdw"""
    spPr = shape._element.spPr
    for e in spPr.findall(qn('a:effectLst')):
        spPr.remove(e)
    lst = spPr.makeelement(qn('a:effectLst'), {})
    sh = lst.makeelement(qn('a:outerShdw'), {
        'blurRad': str(int(blur_px * PXE)), 'dist': str(int(dist_px * PXE)),
        'dir': '5400000', 'rotWithShape': '0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val': '000000'})
    clr.append(clr.makeelement(qn('a:alpha'),
                               {'val': str(int(alpha * 100000))}))
    sh.append(clr)
    lst.append(sh)
    ln = spPr.find(qn('a:ln'))
    (ln.addprevious(lst) if ln is not None else spPr.append(lst))


# ---------- texto ----------
def style_run(r, text, size_px, rgb, bold=False, italic=False, spacing=0.0,
              caps=False):
    r.text = text
    f = r.font
    f.name = FONT
    f.size = pt(size_px)
    f.bold = bold
    f.italic = italic
    if rgb is not None:
        f.color.rgb = rgb
    rPr = r._r.get_or_add_rPr()
    rPr.set('lang', LANG)
    if spacing:
        rPr.set('spc', str(int(round(spacing * 75))))   # 1/100 pt
    if caps:
        rPr.set('cap', 'all')
    return r


def textbox(slide, x, y, w, h, runs, line_px=None, align='left',
            anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(max(h, 1)))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = ALIGN.get(align, PP_ALIGN.LEFT)
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    if line_px:
        p.line_spacing = pt(line_px)
    for r in runs:
        if r.get('br'):
            br = p._p.makeelement(qn('a:br'), {})
            p._p.append(br)
            continue
        if not r.get('t'):
            continue
        style_run(p.add_run(), r['t'], r.get('size', 16), hexcolor(r.get('color')),
                  r.get('bold', False), r.get('italic', False),
                  r.get('spacing', 0.0), r.get('caps', False))
    return tb


def contido(box, arq):
    """object-fit:contain -- a tela cabe inteira na caixa do <img>, centrada"""
    with Image.open(arq) as im:
        nw, nh = im.size
    k = min(box['w'] / nw, box['h'] / nh)
    w, h = nw * k, nh * k
    return box['x'] + (box['w'] - w) / 2, box['y'] + (box['h'] - h) / 2, w, h


# ---------- blocos (fundo + bordas do CSS) ----------
EDGES = (('Top', 't'), ('Right', 'r'), ('Bottom', 'b'), ('Left', 'l'))


def nome(nd, sufixo):
    base = (nd['cls'] or nd['tag']).split()[0]
    return '%s · %s' % (base, sufixo)


def draw_box(slide, nd):
    st, b = nd['style'], nd['box']
    fill, fa = color(st['backgroundColor'])
    sides = {}
    for css, k in EDGES:
        w = num(st['border%sWidth' % css])
        c, a = color(st['border%sColor' % css])
        if w > 0 and st['border%sStyle' % css] != 'none' and a > 0:
            sides[k] = (w, c, a)
    uniform = (len(sides) == 4
               and len({(round(w, 2), str(c), a) for w, c, a in sides.values()}) == 1)
    shp = None
    if uniform:
        w, c, a = sides['t']
        # a borda do PowerPoint fica centrada na aresta: recua meio traco
        shp = rect(slide, b['x'] + w / 2, b['y'] + w / 2, b['w'] - w, b['h'] - w,
                   fill=fill, alpha=fa, line=c, lw=w, line_alpha=a)
    elif fill is not None and fa > 0:
        shp = rect(slide, b['x'], b['y'], b['w'], b['h'], fill=fill, alpha=fa)
    if not uniform:
        for k, (w, c, a) in sides.items():           # bordas assimetricas
            if k == 't':
                r = (b['x'], b['y'], b['w'], w)
            elif k == 'b':
                r = (b['x'], b['y'] + b['h'] - w, b['w'], w)
            elif k == 'l':
                r = (b['x'], b['y'], w, b['h'])
            else:
                r = (b['x'] + b['w'] - w, b['y'], w, b['h'])
            rect(slide, *r, fill=c, alpha=a)
    if shp is not None:
        shp.name = nome(nd, 'fundo')
        if 'moldura' in nd['cls']:
            shadow(shp, 6, 22, 0.10)
    return shp


def draw_text(slide, nd):
    runs = nd.get('runs')
    if not runs:
        return
    st, cb = nd['style'], nd['cbox']
    tb = nd.get('tbox') or cb
    x, w = cb['x'], cb['w']
    if nd['cls'] == 'q':          # item de texto anonimo do flex, depois da aspa
        x = tb['x']
        w = cb['x'] + cb['w'] - tb['x']
    line = num(st['lineHeight'], 0) or None
    # caixa de uma linha so (pill, tag de banco, rotulo): a largura no HTML e
    # justa ao texto, entao a quebra fica desligada para nao sobrar letra
    corpo = max([r.get('size', 16) for r in runs if not r.get('br')] or [16])
    uma_linha = tb['h'] <= (line or corpo * 1.25) * 1.6
    y = topo_caixa(topo_linha(tb['y'], corpo, line), corpo, line)
    tx = textbox(slide, x, y, w, tb['h'], runs, line, st['textAlign'],
                 wrap=not uma_linha)
    tx.name = nome(nd, 'texto')


# ---------- pseudo-elementos ----------
def draw_pseudo(slide, nd, key):
    p = nd.get(key)
    if not p:
        return
    if p.get('display') == 'none':
        return
    b = nd['box']
    fill, fa = color(p['backgroundColor'])
    bw = num(p['borderLeftWidth'])
    txt = (p.get('content') or '').strip('"')
    left, top = num(p['left'], None), num(p['top'], None)
    pw, ph = num(p['width'], 0), num(p['height'], 0)

    if bw > 0 and (fill is None or fa == 0):         # seta do ciclo
        bt = num(p['borderTopWidth'])
        triangle(slide, b['x'] + left, b['y'] + top - bt, bw, bt * 2,
                 hexcolor(p['borderLeftColor']))
        return
    if fill is not None and fa > 0:                  # regua verde, bolinha da lista
        shape = (MSO_SHAPE.OVAL if '%' in (p.get('borderRadius') or '')
                 else MSO_SHAPE.RECTANGLE)
        rect(slide, b['x'] + left, b['y'] + top, pw, ph, fill=fill, alpha=fa,
             shape=shape)
        return
    if not txt:
        return
    corpo = num(p['fontSize'], 16)
    run = [{'t': txt, 'size': corpo, 'color': p['color'],
            'bold': int(num(p['fontWeight'], 400)) >= 600,
            'spacing': num(p['letterSpacing'], 0.0)}]
    line = num(p['lineHeight'], 0) or None
    if nd['tag'] == 'section':                       # numero da pagina
        right = SW - num(p['right'])
        textbox(slide, right - 120, topo_caixa(b['y'] + top, corpo, line), 120,
                ph or 15, run, line, 'right')
    elif p['position'] == 'static':                  # aspa da objecao
        cb = nd['cbox']
        textbox(slide, cb['x'],
                topo_caixa(cb['y'] + num(p['marginTop']), corpo, line),
                pw or 20, ph or 20, run, line, 'left')
    else:                                            # check da lista
        textbox(slide, b['x'] + (left or 0),
                topo_caixa(b['y'] + (top or 0), corpo, line), pw or 20,
                ph or line or 20, run, line, 'left')


# ---------- tabela nativa (pagina "antes x depois") ----------
def draw_table(slide, nodes):
    tnode = next(n for n in nodes if n['tag'] == 'table')
    cells = [n for n in nodes if n['tag'] in ('th', 'td')]
    rows = []
    for c in cells:
        y = round(c['box']['y'], 1)
        if not rows or rows[-1][0] != y:
            rows.append((y, []))
        rows[-1][1].append(c)
    ncol = len(rows[0][1])
    b = tnode['box']
    gf = slide.shapes.add_table(len(rows), ncol, emu(b['x']), emu(b['y']),
                                emu(b['w']), emu(b['h']))
    tbl = gf.table
    set_table_plain(tbl)
    for i, c in enumerate(rows[0][1]):
        tbl.columns[i].width = emu(c['box']['w'])
    set_row_heights(tbl, [r[1][0]['box']['h'] * PXE / 914400 for r in rows])
    for ri, (_, rcells) in enumerate(rows):
        for ci, nd in enumerate(rcells):
            st = nd['style']
            cell = tbl.cell(ri, ci)
            fill, fa = color(st['backgroundColor'])
            if fill is None or fa == 0:
                cell.fill.background()
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill
            for css, edge in (('Top', 'T'), ('Right', 'R'),
                              ('Bottom', 'B'), ('Left', 'L')):
                w = num(st['border%sWidth' % css])
                c, a = color(st['border%sColor' % css])
                cell_border(cell, edge, c if (w > 0 and a > 0) else None,
                            (w * 0.75) or 0.75)
            cell.margin_left = emu(num(st['paddingLeft']))
            cell.margin_right = emu(num(st['paddingRight']))
            cell.margin_top = emu(num(st['paddingTop']))
            cell.margin_bottom = emu(num(st['paddingBottom']))
            cell.vertical_anchor = MSO_ANCHOR.TOP
            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = ALIGN.get(st['textAlign'], PP_ALIGN.LEFT)
            line = num(st['lineHeight'], 0)
            if line:
                p.line_spacing = pt(line)
            if ri:                       # divisoria tambem no topo de baixo
                prev = rows[ri - 1][1][ci]['style']
                w = num(prev['borderBottomWidth'])
                c, a = color(prev['borderBottomColor'])
                if w > 0 and a > 0:
                    cell_border(cell, 'T', c, w * 0.75)
            for r in nd['runs']:
                style_run(p.add_run(), r['t'], r.get('size', 16),
                          hexcolor(r.get('color')), r.get('bold', False),
                          r.get('italic', False), r.get('spacing', 0.0),
                          r.get('caps', False))


# ---------- montagem ----------
def build(layout, telas, out):
    prs = Presentation()
    prs.slide_width = emu(SW)
    prs.slide_height = emu(SH)
    blank = prs.slide_layouts[6]

    for sl in layout['slides']:
        slide = prs.slides.add_slide(blank)
        nodes = sl['nodes']
        tabela = any(n['tag'] == 'table' for n in nodes)
        pulados = {'thead', 'tbody', 'tr', 'th', 'td'} if tabela else set()

        for nd in nodes:
            if nd['tag'] in pulados:
                continue
            if nd['tag'] == 'section':
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                for chave, (ang, stops) in GRADIENTE.items():
                    if chave in nd['cls'].split():
                        grad_rect(slide, 0, 0, SW, SH, ang, stops)
            elif nd['tag'] == 'img':
                arq = telas / nd['img']
                x, y, w, h = contido(nd['box'], arq)
                slide.shapes.add_picture(str(arq), emu(x), emu(y), emu(w),
                                         emu(h))
            else:
                draw_box(slide, nd)
            draw_pseudo(slide, nd, 'before')
            draw_text(slide, nd)
            draw_pseudo(slide, nd, 'after')

        if tabela:
            draw_table(slide, nodes)

    prs.core_properties.title = ('Simule com os bancos — Treinamento para '
                                 'Consultores de Vendas')
    prs.core_properties.author = 'Localiza Seminovos'
    prs.save(out)
    return prs


if __name__ == '__main__':
    raiz = pathlib.Path(__file__).resolve().parent.parent
    lay = pathlib.Path(sys.argv[1]) if len(sys.argv) > 1 else raiz / 'scripts/layout.json'
    telas = pathlib.Path(sys.argv[2]) if len(sys.argv) > 2 else raiz / 'scripts/telas'
    out = pathlib.Path(sys.argv[3]) if len(sys.argv) > 3 else \
        raiz / 'output/treinamento_simulador_pdv.pptx'
    prs = build(json.load(open(lay, encoding='utf-8')), telas, out)
    print('paginas:', len(prs.slides.__iter__.__self__._sldIdLst),
          '->', out)
