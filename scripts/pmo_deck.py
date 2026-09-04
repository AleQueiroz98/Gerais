# -*- coding: utf-8 -*-
"""Deck de 3 paginas com o status das quatro frentes de F&I.

    pagina 1  identidade por frente com o color code do plano de escalada
              (versao da pagina HTML): faixa colorida, responsavel, valor
              gerado e status em pastilha com rotulo escrito
    pagina 2  a pagina do print, gerada por `pmo_status.py` sem alteracao
    pagina 3  combinacao das duas: color code e identidade da frente da 1,
              titulo a partir das contagens e contagem por frente da 2, com a
              coluna "Pontos a escalar" promovida a faixa de decisoes pedidas

Tudo em tabelas nativas do PowerPoint, editaveis celula a celula. As barras de
cor sao a primeira coluna da tabela (e nao shapes soltos) para continuarem
coladas na linha quando o texto for editado.

Conteudo: pmo_deck_content.py (paginas 1 e 3) e pmo_content.py (pagina 2)
Tokens e helpers: deckstyle.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from deckstyle import (
    BLACK, TEXT_DARK, TEXT_REG, TEXT_HELPER, WHITE,
    GRAY_X_LIGHT, GRAY_LIGHT, GRAY_MEDIUM, PLACEHOLDER,
    FONT_NAME, CELL_MX,
    text_w, nlines, lh, split_bold,
    style_run, cell_box, set_table_plain, set_row_heights,
)
from pmo_status import (
    SW, SH, SEM, GLIFO, DOTS, rgb, fh, fix_para, blank, ctext, textbox, rect,
)
import pmo_status
import pmo_deck_content as C

OUT = '../output/260903__Status_frentes_3paginas.pptx'

ML = 0.42
CW = SW - 2 * ML
TOP = 0.30
FOOT_Y = SH - 0.34
PAD_ROW = 0.13
CELL_MY = 0.04

# color code das frentes, do plano de escalada de F&I; o ocre foi escurecido
# um passo para o texto claro sobre ele continuar legivel na projecao
FRENTE_RGB = [RGBColor(0x40, 0x70, 0x8F), RGBColor(0x8C, 0x3A, 0x6B),
              RGBColor(0x9C, 0x82, 0x26), RGBColor(0x59, 0x59, 0x59)]
CHIP_BG = RGBColor(0x34, 0x34, 0x32)
RED = RGBColor(0xCC, 0x00, 0x00)

PT_TITLE, PT_SUB, PT_LEGEND, PT_HEAD, PT_FOOT = 15.0, 11.5, 8.0, 8.0, 8.0
PT_FRENTE, PT_META = 10.0, 8.0
BODY_STEPS = (9.5, 9.0, 8.5, 8.0, 7.5)

SUB_LABEL = {'ok': 'no prazo', 'plano': 'dentro do plano',
             'risco': 'prazo a definir', 'atraso': 'fora do plano'}
SHORT_LABEL = {'ok': 'Concluído', 'plano': 'Em andamento',
               'risco': 'Bloqueado', 'atraso': 'Em andamento'}


# --------------------------------------------------------------- helpers
def tb_runs(slide, x, y, w, h, parts, size, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, wrap=True):
    """caixa de texto com varios runs: parts = [(texto, negrito, cor)]"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for seg, bold, color in parts:
        style_run(p.add_run(), seg, size, color, bold)
    fix_para(tf, size)
    return tb


def cell_paras(cell, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
               mx=CELL_MX, mt=CELL_MY, mb=CELL_MY, space=2.0):
    """celula com varios paragrafos.

    paras = [(size, [(texto, bold, cor) | (texto, bold, cor, size)])]
    O quarto item permite um run com tamanho proprio no mesmo paragrafo, que e
    como o simbolo do semaforo fica na mesma linha do rotulo.
    """
    cell.margin_left = Inches(mx); cell.margin_right = Inches(mx)
    cell.margin_top = Inches(mt); cell.margin_bottom = Inches(mb)
    cell.vertical_anchor = anchor
    tf = cell.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, (size, parts) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if i:
            p.space_before = Pt(space)
        for part in parts:
            seg, bold, color = part[0], part[1], part[2]
            style_run(p.add_run(), seg, part[3] if len(part) > 3 else size,
                      color, bold)
        p.font.size = Pt(size)
        p.font.name = FONT_NAME
    return tf


def glyph(key, size):
    """(caractere, tamanho) do semaforo: check para concluido, tres pontos nos demais"""
    return (GLIFO['ok'], size + 3.0) if key == 'ok' else (DOTS, size + 1.5)


def legend_row(slide, x, y, w, itens, size=PT_LEGEND, counts=False):
    """legenda em uma linha, a partir de x, com contagem opcional"""
    gw = 0.26
    xx = x
    for it in itens:
        key, label = it[0], it[1]
        ch, gs = glyph(key, size)
        textbox(slide, xx, y - 0.05, gw, 0.24, ch, size=gs,
                color=rgb(SEM[key][1]), bold=True, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        lw = text_w(label, size)
        textbox(slide, xx + gw + 0.04, y, lw + 0.06, 0.16, label,
                size=size, color=TEXT_REG, wrap=False)
        xx += gw + 0.04 + lw + 0.10
        if counts:
            n = str(it[2])
            nw = text_w(n, size, True)
            textbox(slide, xx, y, nw + 0.06, 0.16, n, size=size,
                    color=RED if key == 'atraso' else BLACK, bold=True,
                    wrap=False)
            xx += nw + 0.06
        xx += 0.22
    return xx


def norm_cols(weights):
    s = CW / sum(weights)
    return [w * s for w in weights]


def solve_rows(frentes, avail, row_min, band_h_of, head_h):
    """escolhe o corpo e as alturas de linha que caibam em uma pagina.

    row_min(marco, body_pt): altura minima da linha, calculada pela pagina a
        partir das suas proprias colunas
    band_h_of(frente, body_pt): altura minima da faixa da frente (celula
        mesclada verticalmente), distribuida entre as linhas do grupo
    """
    for body_pt in BODY_STEPS:
        heights = [head_h]
        groups = []
        for f in frentes:
            ix = []
            for m in f['marcos']:
                ix.append(len(heights))
                heights.append(max(0.30, row_min(m, body_pt)))
            groups.append((f, ix))
        # a faixa da frente nao pode ficar mais alta que o grupo
        for f, ix in groups:
            need = band_h_of(f, body_pt)
            have = sum(heights[i] for i in ix)
            if need > have:
                extra = (need - have) / len(ix)
                for i in ix:
                    heights[i] += extra
        if sum(heights) <= avail:
            break

    item_ix = [i for _, ix in groups for i in ix]
    caps = {i: max(heights[i], 0.72) for i in item_ix}
    for _ in range(80):
        slack = avail - sum(heights)
        movable = [i for i in item_ix if heights[i] < caps[i] - 1e-4]
        if slack <= 0.01 or not movable:
            break
        step = slack / len(movable)
        for i in movable:
            heights[i] = min(caps[i], heights[i] + step)
    if sum(heights) > avail + 0.02:
        raise SystemExit('a tabela nao cabe na pagina: %.2f" de conteudo em '
                         '%.2f" disponiveis, ja no corpo minimo de %.1fpt'
                         % (sum(heights), avail, body_pt))
    return body_pt, heights


def body_cell(tbl, ri, ci, val, body_pt, bold_map, sep, left=None):
    """celula de texto do corpo, com NA/TBD em cinza claro e negrito seletivo"""
    c = tbl.cell(ri, ci)
    cell_box(c, None, (left, None, None, sep), lw=0.5)
    vazio = val in (C.NA, C.TBD)
    ctext(c, [(val, False)] if vazio else split_bold(val, bold_map.get(val, [])),
          body_pt, PLACEHOLDER if vazio else TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)
    return c


def feito_cell(tbl, ri, ci, m, body_pt, with_prazo):
    """o que precisava ser feito: titulo em negrito, complemento em cinza e,
    quando `with_prazo`, o prazo e a nota logo abaixo"""
    c = tbl.cell(ri, ci)
    parts = [(m['feito'], True, BLACK)]
    if m.get('qual'):
        parts.append((m['qual'], False, TEXT_REG))
    paras = [(body_pt, parts)]
    if with_prazo:
        linha = [(m['prazo'], False,
                  PLACEHOLDER if m['prazo'] == C.TBD else TEXT_DARK)]
        if m.get('nota'):
            linha.append(('   ' + m['nota'], False, RED))
        paras.append((body_pt - 0.5, linha))
    cell_paras(c, paras, anchor=MSO_ANCHOR.MIDDLE)
    return c


def status_pill(tbl, ri, ci, key, body_pt):
    """pastilha de status: fundo tingido, simbolo na linha do rotulo e o
    qualificador logo abaixo, como na pagina HTML"""
    _, fg, bg = SEM[key]
    c = tbl.cell(ri, ci)
    cell_box(c, rgb(bg), (WHITE, WHITE, WHITE, WHITE), lw=3.0)
    ch, gs = glyph(key, body_pt - 1.0)
    cell_paras(c, [
        (body_pt - 1.0, [(ch + ' ', True, rgb(fg), gs),
                         (SHORT_LABEL[key], True, BLACK)]),
        (body_pt - 1.5, [(SUB_LABEL[key], False, TEXT_DARK)]),
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, mx=0.03, space=1.0)
    return c


def status_h(body_pt):
    """altura minima da pastilha de status"""
    return fh(body_pt + 2.0) + fh(body_pt - 1.5) + 1.0 / 72.0 + 2 * CELL_MY


def footer(slide, fonte, classe, page):
    textbox(slide, ML, FOOT_Y, CW - 3.4, 0.20, fonte, size=PT_FOOT,
            color=TEXT_HELPER)
    textbox(slide, ML + CW - 3.4, FOOT_Y, 3.15, 0.20, classe, size=PT_FOOT,
            color=TEXT_HELPER, align=PP_ALIGN.RIGHT)
    textbox(slide, ML + CW - 0.20, FOOT_Y, 0.20, 0.20, str(page), size=PT_FOOT,
            color=TEXT_DARK, bold=True, align=PP_ALIGN.RIGHT)


def chip(slide, x, y, label):
    """tag escura com o canto direito chanfrado, como a do plano de escalada"""
    w = text_w(label, PT_LEGEND, True) + 0.46
    s = rect(slide, x, y, w, 0.26, fill=CHIP_BG, shape=MSO_SHAPE.PENTAGON,
             adj=0.10)
    tf = s.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style_run(p.add_run(), label, PT_LEGEND, WHITE, True)
    fix_para(tf, PT_LEGEND)
    return w


def band_sizes(body_pt):
    """tamanhos da faixa da frente, derivados do corpo: quando o corpo reduz
    para caber na pagina, a faixa reduz junto"""
    return body_pt + 1.5, max(6.5, body_pt - 1.0)


def band_paras(fi, f, col, pt_nome, pt_meta, tally=None, meta_inline=False):
    """paragrafos da faixa: numero e nome, responsavel/valor e contagem"""
    paras = [(pt_nome, [('%d. ' % (fi + 1), True, col),
                        (f['nome'], True, BLACK)])]
    if meta_inline:
        paras.append((pt_meta, [(f['resp'], False, TEXT_DARK),
                                ('  ' + chr(0x2022) + '  ', False, GRAY_MEDIUM),
                                (f['valor'], True, BLACK)]))
    else:
        paras.append((pt_meta, [('RESP.  ', True, TEXT_HELPER),
                                (f['resp'], False, TEXT_DARK)]))
        paras.append((pt_meta, [('VALOR  ', True, TEXT_HELPER),
                                (f['valor'], True, BLACK)]))
    if tally:
        paras.append((pt_meta, tally))
    return paras


def band_h(f, w, pt_nome, pt_meta, tally=False, meta_inline=False):
    """altura minima da faixa da frente, para o grupo nunca ficar mais curto"""
    nl = nlines('%d. ' % 1 + f['nome'], pt_nome, w, True)
    if meta_inline:
        ml = nlines(f['resp'] + '  ' + chr(0x2022) + '  ' + f['valor'],
                    pt_meta, w)
    else:
        ml = nlines('RESP.  ' + f['resp'], pt_meta, w) + 1
    return nl * fh(pt_nome) + (ml + (1 if tally else 0)) * fh(pt_meta) + 0.32


def frente_band(tbl, ri, n, fi, f, paras):
    """barra colorida da frente + celula de identidade, ambas mescladas
    verticalmente sobre as linhas do grupo"""
    col = FRENTE_RGB[fi % len(FRENTE_RGB)]
    if n > 1:
        tbl.cell(ri, 0).merge(tbl.cell(ri + n - 1, 0))
        tbl.cell(ri, 1).merge(tbl.cell(ri + n - 1, 1))
    blank(tbl.cell(ri, 0), col)
    c = tbl.cell(ri, 1)
    cell_box(c, GRAY_X_LIGHT, (None, None, None, GRAY_MEDIUM), lw=0.5)
    cell_paras(c, paras, mx=0.12, mt=0.10, space=2.5)
    return col


def head_row(tbl, names, center_ix):
    for ci, name in enumerate(names):
        c = tbl.cell(0, ci)
        bd = (None, None, None, BLACK)
        align = PP_ALIGN.CENTER if ci == center_ix else PP_ALIGN.LEFT
        if name:
            cell_box(c, GRAY_X_LIGHT, bd, lw=1.0)
            ctext(c, [(name.upper(), True)], PT_HEAD, BLACK, align=align,
                  anchor=MSO_ANCHOR.MIDDLE)
        else:
            blank(c, GRAY_X_LIGHT, bd, lw=1.0)


def new_table(slide, y, widths, nrows, heights):
    gf = slide.shapes.add_table(nrows, len(widths), Inches(ML), Inches(y),
                                Inches(CW), Inches(sum(heights)))
    tbl = gf.table
    set_table_plain(tbl)
    set_row_heights(tbl, heights)
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)
    return tbl


def title_block(slide, y, bold, tail):
    tit_n = nlines(bold + tail, PT_TITLE, CW, True)
    tb_runs(slide, ML, y, CW, tit_n * lh(PT_TITLE),
            [(bold, True, BLACK), (tail, False, BLACK)], PT_TITLE)
    y += tit_n * lh(PT_TITLE) + 0.10
    rect(slide, ML, y, CW, 0.024, fill=RED)
    return y + 0.024


# --------------------------------------------------------------- pagina 1
P1_W = norm_cols([0.09, 1.80, 2.48, 1.25, 2.38, 2.38, 2.12])
P1_HEAD = ['', 'Frente', 'O que precisava ser feito', 'Status',
           'Progresso recente', 'Próximos passos', 'Pontos a escalar']
P1_TEXT = ((4, 'prog'), (5, 'prox'), (6, 'esc'))


def p1_row_h(m, body_pt):
    """a coluna do entregavel tem duas linhas de texto (titulo e prazo), e a
    pastilha de status tem altura propria; a linha usa a maior das tres"""
    titulo = m['feito'] + (m.get('qual') or '')
    h_feito = (nlines(titulo, body_pt, P1_W[2] - 2 * CELL_MX, True) * fh(body_pt)
               + fh(body_pt - 0.5) + 2 * CELL_MY + 0.04)
    h_text = max(nlines(m[k], body_pt, P1_W[ci] - 2 * CELL_MX)
                 for ci, k in P1_TEXT) * fh(body_pt) + PAD_ROW
    return max(h_feito, status_h(body_pt), h_text)


def p1_band_h(f, body_pt):
    pn, pm = band_sizes(body_pt)
    return band_h(f, P1_W[1] - 0.24, pn, pm)


def page1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, SW, SH, fill=WHITE)

    y = title_block(slide, TOP, C.P1_TITULO_BOLD, C.P1_TITULO) + 0.13
    cw = chip(slide, ML, y - 0.04, C.P1_CHIP)
    legend_row(slide, ML + cw + 0.28, y, CW, C.P1_LEGENDA, counts=True)
    textbox(slide, ML + CW - 1.30, y, 1.30, 0.18, C.P1_MARCA, size=PT_LEGEND,
            color=BLACK, bold=True, align=PP_ALIGN.RIGHT, wrap=False)
    y += 0.32

    head_h = fh(PT_HEAD) + 0.13
    body_pt, heights = solve_rows(C.P1, FOOT_Y - 0.14 - y, p1_row_h,
                                  p1_band_h, head_h)
    nrows = 1 + sum(len(f['marcos']) for f in C.P1)
    tbl = new_table(slide, y, P1_W, nrows, heights)
    head_row(tbl, P1_HEAD, 3)

    ri = 1
    for fi, f in enumerate(C.P1):
        n = len(f['marcos'])
        col = FRENTE_RGB[fi % len(FRENTE_RGB)]
        pn, pm = band_sizes(body_pt)
        frente_band(tbl, ri, n, fi, f, band_paras(fi, f, col, pn, pm))
        for k, m in enumerate(f['marcos']):
            sep = GRAY_MEDIUM if k == n - 1 else GRAY_LIGHT
            feito_cell(tbl, ri, 2, m, body_pt, with_prazo=True)
            cell_box(tbl.cell(ri, 2), None, (None, None, None, sep), lw=0.5)
            status_pill(tbl, ri, 3, m['status'], body_pt)
            for ci, key in ((4, 'prog'), (5, 'prox'), (6, 'esc')):
                body_cell(tbl, ri, ci, m[key], body_pt, C.P1_BOLD, sep)
            ri += 1

    footer(slide, C.P1_FONTE, C.P1_CLASSE, 1)
    return body_pt, y + sum(heights)


# --------------------------------------------------------------- pagina 3
P3_W = norm_cols([0.09, 2.05, 2.60, 0.70, 1.22, 2.80, 2.85])
P3_HEAD = ['', 'Frente', 'O que precisava ser feito', 'Prazo', 'Status',
           'Progresso recente', 'Próximos passos']
P3_TEXT = ((5, 'prog'), (6, 'prox'))
DEC_W = norm_cols([0.09, 2.60, 7.10, 2.60])
DEC_ROW = 0.26


def p3_row_h(m, body_pt):
    h_feito = (nlines(m['feito'], body_pt, P3_W[2] - 2 * CELL_MX, True)
               * fh(body_pt) + PAD_ROW)
    h_prazo = fh(body_pt) + (fh(body_pt - 1.0) if m.get('nota') else 0) + PAD_ROW
    h_text = max(nlines(m[k], body_pt, P3_W[ci] - 2 * CELL_MX)
                 for ci, k in P3_TEXT) * fh(body_pt) + PAD_ROW
    return max(h_feito, h_prazo, status_h(body_pt), h_text)


def p3_band_h(f, body_pt):
    pn, pm = band_sizes(body_pt)
    return band_h(f, P3_W[1] - 0.24, pn, pm, tally=True, meta_inline=True)


def p3_resumo(f):
    """contagem por status da frente, na ordem do semaforo"""
    linha = []
    for k in pmo_status.ORDEM:
        n = sum(1 for m in f['marcos'] if m['status'] == k)
        if not n:
            continue
        if linha:
            linha.append(('  ' + chr(0x2022) + '  ', False, GRAY_MEDIUM))
        linha.append(('%d %s' % (n, C.P3_RESUMO[k]), True, rgb(SEM[k][1])))
    return linha


def page3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, SW, SH, fill=WHITE)

    y = title_block(slide, TOP, C.P3_TITULO_BOLD, C.P3_TITULO) + 0.11
    textbox(slide, ML, y, 4.6, lh(PT_SUB), C.P3_SUB, size=PT_SUB,
            color=BLACK, bold=True)
    legend_row(slide, ML + 4.75, y + 0.045, CW - 4.75, C.P3_LEGENDA)
    y += lh(PT_SUB) + 0.11

    # a faixa de decisoes ocupa o pe da pagina; a tabela usa o que sobra
    dec_h = 0.26 + DEC_ROW * len(C.P3_DECISOES)
    dec_y = FOOT_Y - 0.14 - dec_h
    head_h = fh(PT_HEAD) + 0.13
    body_pt, heights = solve_rows(C.P3, dec_y - 0.20 - y, p3_row_h,
                                  p3_band_h, head_h)
    nrows = 1 + sum(len(f['marcos']) for f in C.P3)
    tbl = new_table(slide, y, P3_W, nrows, heights)
    head_row(tbl, P3_HEAD, 4)

    ri = 1
    for fi, f in enumerate(C.P3):
        n = len(f['marcos'])
        col = FRENTE_RGB[fi % len(FRENTE_RGB)]
        pn, pm = band_sizes(body_pt)
        frente_band(tbl, ri, n, fi, f,
                    band_paras(fi, f, col, pn, pm, tally=p3_resumo(f),
                               meta_inline=True))
        for k, m in enumerate(f['marcos']):
            sep = GRAY_MEDIUM if k == n - 1 else GRAY_LIGHT
            feito_cell(tbl, ri, 2, m, body_pt, with_prazo=False)
            cell_box(tbl.cell(ri, 2), None, (None, None, None, sep), lw=0.5)
            c = tbl.cell(ri, 3)
            cell_box(c, None, (None, None, None, sep), lw=0.5)
            paras = [(body_pt, [(m['prazo'], False,
                                 PLACEHOLDER if m['prazo'] == C.TBD
                                 else TEXT_DARK)])]
            if m.get('nota'):
                paras.append((body_pt - 1.0, [(m['nota'], False, RED)]))
            cell_paras(c, paras, anchor=MSO_ANCHOR.MIDDLE)
            status_pill(tbl, ri, 4, m['status'], body_pt)
            for ci, key in ((5, 'prog'), (6, 'prox')):
                body_cell(tbl, ri, ci, m[key], body_pt, C.P3_BOLD, sep)
            ri += 1

    # ---------- faixa de decisoes pedidas ----------
    tb_runs(slide, ML, dec_y, CW, 0.24,
            [(C.P3_DECISOES_TITULO, True, BLACK),
             ('   as três frentes que dependem de uma definição desta reunião',
              False, TEXT_REG)], PT_SUB - 1.5)
    ty = dec_y + 0.26
    nomes = [f['nome'] for f in C.P3]
    dtbl = new_table(slide, ty, DEC_W, len(C.P3_DECISOES),
                     [DEC_ROW] * len(C.P3_DECISOES))
    for i, (frente, dec, resp) in enumerate(C.P3_DECISOES):
        blank(dtbl.cell(i, 0), FRENTE_RGB[nomes.index(frente) % len(FRENTE_RGB)])
        for ci, (val, bold, color) in enumerate(
                ((frente, True, BLACK), (dec, False, TEXT_DARK),
                 (resp, False, TEXT_REG)), start=1):
            c = dtbl.cell(i, ci)
            cell_box(c, GRAY_X_LIGHT, (None, None, None, WHITE), lw=1.5)
            ctext(c, [(val, bold)], PT_META, color,
                  align=PP_ALIGN.RIGHT if ci == 3 else PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.MIDDLE, mx=0.12)

    footer(slide, C.P3_FONTE, C.P3_CLASSE, 3)
    return body_pt, ty + DEC_ROW * len(C.P3_DECISOES)


# --------------------------------------------------------------- main
def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
    pt1, base1 = page1(prs)
    info2 = pmo_status.build(prs)
    pt3, base3 = page3(prs)
    prs.core_properties.title = 'Status das quatro frentes de F&I'
    prs.save(OUT)
    print('deck de 3 paginas | corpo p1 %.1fpt (base %.2f") | '
          'p2 %.1fpt (base %.2f") | p3 %.1fpt (base %.2f")'
          % (pt1, base1, info2['corpo_pt'], info2['base'], pt3, base3))


if __name__ == '__main__':
    main()
