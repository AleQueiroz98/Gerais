# -*- coding: utf-8 -*-
"""Pagina unica de status das quatro frentes (PMO), na identidade do painel.

Mesma linguagem visual do deck `Acompanhamento das frentes`: barra verde,
tabela nativa do PowerPoint, cabecalhos em cinza claro e corpo que reduz de
9,5pt ate 7,5pt para caber em uma pagina.

O semaforo e desenhado com vetores (check e tres pontos), nao com fontes de
simbolo, para nao depender de Wingdings/Segoe na maquina de quem abre.

Conteudo: pmo_content.py    Tokens e helpers: deckstyle.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from deckstyle import (
    GREEN, GREEN_DEEP, BLACK, TEXT_DARK, TEXT_REG, TEXT_HELPER,
    WHITE, GRAY_X_LIGHT, GRAY_LIGHT, GRAY_MEDIUM, PLACEHOLDER,
    FONT_NAME, CELL_MX,
    text_w, nlines, lh, split_bold,
    no_shadow, style_run, cell_box, cell_text, set_table_plain, set_row_heights,
)
import pmo_content as C

OUT = '../output/260903__Status_quatro_frentes.pptx'

# ---------- geometria ----------
SW, SH = 13.3333, 7.5
ML = 0.42
CW = SW - 2 * ML
TOP = 0.30
BOTTOM = 7.05              # base da tabela; o rodape fica abaixo
PAD_ROW = 0.13             # folga vertical da linha, acima das margens da celula
CELL_MY = 0.04             # margem vertical dentro da celula

PT_TITLE, PT_SECTION, PT_LEGEND, PT_HEAD, PT_FOOT = 15.0, 11.5, 8.0, 8.0, 8.0
PT_FRENTE = 10.0
BODY_STEPS = (9.5, 9.0, 8.5, 8.0, 7.5)

# ---------- semaforo ----------
# chave -> (rotulo, cor do simbolo, fundo da pastilha)
SEM = {
    'ok':     ('Concluído',                     (0x14, 0x6A, 0x3D), (0xE3, 0xF5, 0xEA)),
    'plano':  ('Em andamento, dentro do plano', (0x4C, 0x8E, 0x14), (0xEF, 0xF8, 0xE4)),
    'risco':  ('Em andamento, com riscos',      (0xC2, 0x8E, 0x00), (0xFF, 0xF6, 0xDC)),
    'atraso': ('Em andamento, fora do plano',   (0xC0, 0x1F, 0x1F), (0xFB, 0xE7, 0xE7)),
}
ORDEM = ('ok', 'plano', 'risco', 'atraso')
RESUMO = {'ok': 'concluído', 'plano': 'no plano',
          'risco': 'com riscos', 'atraso': 'fora do plano'}

GLIFO = {'ok': '\u2713'}          # check; as demais usam tres pontos
DOTS = '\u2022\u2009\u2022\u2009\u2022'

COLS = [
    ('',                          0.085, PP_ALIGN.LEFT),
    ('O que precisava ser feito', 2.78,  PP_ALIGN.LEFT),
    ('Prazo',                     0.76,  PP_ALIGN.LEFT),
    ('Status',                    0.66,  PP_ALIGN.CENTER),
    ('Progresso recente',         2.71,  PP_ALIGN.LEFT),
    ('Próximos passos',           2.71,  PP_ALIGN.LEFT),
    ('Pontos a escalar',          2.71,  PP_ALIGN.LEFT),
]
_s = CW / sum(c[1] for c in COLS)
COLS = [(n, w * _s, a) for n, w, a in COLS]
WIDTHS = [c[1] for c in COLS]


def rgb(t):
    return RGBColor(*t)


def fh(pt):
    """altura de linha usada no calculo das alturas de linha da tabela.

    Um pouco maior que lh(): o PowerPoint trata a altura da linha como minimo,
    e uma estimativa curta faz a tabela crescer e passar do rodape.
    """
    return pt * 1.28 / 72.0


def fix_para(tf, size):
    """fixa o tamanho no paragrafo, nao so no run.

    Sem isso o paragrafo mantem o padrao de 18pt no endParaRPr e o
    PowerPoint/LibreOffice usa esse valor como altura minima da linha,
    inflando a tabela.
    """
    for p in tf.paragraphs:
        p.font.size = Pt(size)
        p.font.name = FONT_NAME
    return tf


def resumo(cell, marcos, size):
    """contagem por status, colorida, no canto direito da faixa da frente"""
    ordem = [k for k in ORDEM if any(m[2] == k for m in marcos)]
    tf = cell.text_frame
    tf.word_wrap = False
    tf.clear()
    cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.16)
    cell.margin_top = cell.margin_bottom = Inches(CELL_MY)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    for i, k in enumerate(ordem):
        n = sum(1 for m in marcos if m[2] == k)
        if i:
            style_run(p.add_run(), '   ' + chr(0x2022) + '   ', size, GRAY_MEDIUM)
        rot = RESUMO[k] + ('s' if n > 1 and k == 'ok' else '')
        style_run(p.add_run(), '%d %s' % (n, rot), size, rgb(SEM[k][1]), bold=True)
    return fix_para(tf, size)


def blank(cell, fill=None, borders=(None, None, None, None), lw=0.75):
    """celula sem texto: zera margens e o paragrafo vazio.

    Um paragrafo vazio herda 18pt e o PowerPoint usa isso como altura minima
    da linha (~0,4"), inflando a barra verde e o cabecalho.
    """
    cell_box(cell, fill, borders, lw)
    cell.margin_top = cell.margin_bottom = Inches(0)
    cell.margin_left = cell.margin_right = Inches(0)
    fix_para(cell.text_frame, 1)
    return cell


def ctext(cell, parts, size, color, **kw):
    kw.setdefault('mt', CELL_MY)
    kw.setdefault('mb', CELL_MY)
    return fix_para(cell_text(cell, parts, size, color, **kw), size)


def textbox(slide, x, y, w, h, text, size=9, color=TEXT_DARK, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    style_run(p.add_run(), text, size, color, bold)
    fix_para(tf, size)
    return tb


def rect(slide, x, y, w, h, fill=None, line=None, lw=0.75,
         shape=MSO_SHAPE.RECTANGLE, adj=None):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    no_shadow(s)
    if adj is not None:
        try: s.adjustments[0] = adj
        except Exception: pass
    tf = s.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return s


def status_cell(cell, key, pt, sep):
    """celula de status: pastilha colorida com check ou tres pontos.

    O simbolo fica dentro da celula (e nao como shape solto por cima) para
    continuar alinhado se a linha crescer ao editar o texto no PowerPoint.
    """
    _, fg, bg = SEM[key]
    cell_box(cell, rgb(bg), (WHITE, WHITE, WHITE, WHITE), lw=3.0)
    ch, size = (GLIFO['ok'], pt + 3.0) if key == 'ok' else (DOTS, pt + 1.5)
    ctext(cell, [(ch, True)], size, rgb(fg), align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE, mx=0.02, mt=0.02, mb=0.02)


def legenda(slide, x, y, w):
    """legenda do semaforo, alinhada a direita e em uma unica linha"""
    gw = 0.26
    widths = [gw + 0.06 + text_w(SEM[k][0], PT_LEGEND) + 0.10 for k in ORDEM]
    gap = 0.24
    xx = x + w - (sum(widths) + gap * (len(ORDEM) - 1))
    for k, wi in zip(ORDEM, widths):
        ch, size = (GLIFO['ok'], PT_LEGEND + 3.0) if k == 'ok' \
            else (DOTS, PT_LEGEND + 1.5)
        textbox(slide, xx, y - 0.05, gw, 0.24, ch, size=size,
                color=rgb(SEM[k][1]), bold=True, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        textbox(slide, xx + gw + 0.06, y, wi - gw - 0.06, 0.16, SEM[k][0],
                size=PT_LEGEND, color=TEXT_REG, wrap=False)
        xx += wi + gap


def row_heights(avail):
    """escolhe o corpo e as alturas que caibam em uma pagina"""
    head_h = fh(PT_HEAD) + 0.13
    frente_h = fh(PT_FRENTE) + 0.10
    text_cols = (1, 4, 5, 6)
    for body_pt in BODY_STEPS:
        heights, item_ix = [head_h], []
        for _, marcos in C.FRENTES:
            heights.append(frente_h)
            for m in marcos:
                vals = (m[0], m[3], m[4], m[5])
                n = max(nlines(v, body_pt, WIDTHS[c] - 2 * CELL_MX, c == 1)
                        for c, v in zip(text_cols, vals))
                n = max(n, nlines(m[1], body_pt, WIDTHS[2] - 2 * CELL_MX))
                item_ix.append(len(heights))
                heights.append(max(0.30, n * fh(body_pt) + PAD_ROW))
        if sum(heights) <= avail:
            break

    # sobra distribuida entre as linhas de marco, com teto para nao inchar
    caps = {i: max(heights[i], 0.62) for i in item_ix}
    for _ in range(80):
        slack = avail - sum(heights)
        movable = [i for i in item_ix if heights[i] < caps[i] - 1e-4]
        if slack <= 0.01 or not movable:
            break
        step = slack / len(movable)
        for i in movable:
            heights[i] = min(caps[i], heights[i] + step)
    return body_pt, heights


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, SW, SH, fill=WHITE)

    # ---------- titulo ----------
    marcos = [m for _, ms in C.FRENTES for m in ms]
    cont = {k: sum(1 for m in marcos if m[2] == k) for k in ORDEM}
    titulo = C.TITULO.format(tot=len(marcos), **cont)

    y = TOP
    tit_n = nlines(titulo, PT_TITLE, CW, True)
    textbox(slide, ML, y, CW, tit_n * lh(PT_TITLE), titulo,
            size=PT_TITLE, color=BLACK, bold=True)
    y += tit_n * lh(PT_TITLE) + 0.17

    # ---------- rotulo da secao + legenda ----------
    textbox(slide, ML, y, 4.0, lh(PT_SECTION), C.SUBTITULO,
            size=PT_SECTION, color=BLACK, bold=True)
    legenda(slide, ML + 4.2, y + 0.035, CW - 4.2)
    y += lh(PT_SECTION) + 0.14

    # ---------- tabela ----------
    nrows = 1 + sum(1 + len(ms) for _, ms in C.FRENTES)
    body_pt, heights = row_heights(BOTTOM - y)
    top = y

    gf = slide.shapes.add_table(nrows, len(COLS), Inches(ML), Inches(top),
                                Inches(CW), Inches(sum(heights)))
    tbl = gf.table
    set_table_plain(tbl)
    set_row_heights(tbl, heights)
    for i, w in enumerate(WIDTHS):
        tbl.columns[i].width = Inches(w)

    # cabecalho
    for ci, (name, _, align) in enumerate(COLS):
        c = tbl.cell(0, ci)
        bd = (None, None, None, GRAY_MEDIUM)
        if name:
            cell_box(c, GRAY_X_LIGHT, bd, lw=0.75)
            ctext(c, [(name.upper(), True)], PT_HEAD, TEXT_HELPER,
                  align=align, anchor=MSO_ANCHOR.MIDDLE)
        else:
            blank(c, GRAY_X_LIGHT, bd)

    # corpo
    ri = 1
    for fi, (nome, marcos) in enumerate(C.FRENTES, start=1):
        # faixa da frente: barra verde + titulo numerado (celulas mescladas)
        blank(tbl.cell(ri, 0), GREEN)
        tbl.cell(ri, 1).merge(tbl.cell(ri, 3))
        tbl.cell(ri, 4).merge(tbl.cell(ri, len(COLS) - 1))
        c = tbl.cell(ri, 1)
        cell_box(c, GRAY_LIGHT)
        tf = ctext(c, [('%d. ' % fi, True), (nome, True)], PT_FRENTE, BLACK,
                   anchor=MSO_ANCHOR.MIDDLE, mx=0.14)
        tf.paragraphs[0].runs[0].font.color.rgb = GREEN_DEEP
        c = tbl.cell(ri, 4)
        cell_box(c, GRAY_LIGHT)
        resumo(c, marcos, PT_LEGEND)
        ri += 1

        for k, m in enumerate(marcos):
            feito, prazo, status, prog, prox, esc = m
            sep = GRAY_MEDIUM if k == len(marcos) - 1 else GRAY_LIGHT
            base = (None, None, None, sep)

            blank(tbl.cell(ri, 0), GRAY_X_LIGHT, base, lw=0.5)
            for ci, val, color, bold in ((1, feito, TEXT_DARK, True),
                                         (2, prazo, TEXT_REG, False)):
                c = tbl.cell(ri, ci)
                cell_box(c, None, base, lw=0.5)
                ctext(c, [(val, bold)], body_pt, color,
                      anchor=MSO_ANCHOR.MIDDLE)

            status_cell(tbl.cell(ri, 3), status, body_pt, sep)

            for ci, val in ((4, prog), (5, prox), (6, esc)):
                c = tbl.cell(ri, ci)
                cell_box(c, None, (GRAY_LIGHT, None, None, sep), lw=0.5)
                vazio = val in (C.NA, C.TBD)
                ctext(c, [(val, False)] if vazio
                      else split_bold(val, C.BOLD.get(val, [])),
                      body_pt, PLACEHOLDER if vazio else TEXT_DARK,
                      anchor=MSO_ANCHOR.MIDDLE)
            ri += 1

    # ---------- rodape ----------
    textbox(slide, ML, SH - 0.34, CW, 0.20, C.RODAPE,
            size=PT_FOOT, color=TEXT_HELPER)

    prs.save(OUT)
    print('pagina gerada | titulo %d linha(s) | corpo %.1fpt | %d linhas | '
          'tabela de %.2f" a %.2f"' % (tit_n, body_pt, nrows, top,
                                       top + sum(heights)))


if __name__ == '__main__':
    main()
