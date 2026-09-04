# -*- coding: utf-8 -*-
"""Preview HTML de um .pptx, lido da geometria real do arquivo.

Serve para conferir o layout sem depender do PowerPoint ou do LibreOffice:
percorre shapes e tabelas, le posicao, tamanho, preenchimento e runs de texto
e desenha tudo em divs posicionados (1" = 100px).

    python3 pptx_preview.py ../output/deck.pptx ../output/deck_preview.html

E uma aproximacao: fontes e quebra de linha ficam a cargo do navegador, entao
serve para checar estrutura, alturas e transbordo, nao para prova de cor.
"""
import html
import sys

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu

PX = 100.0   # pixels por polegada


def inches(v):
    return Emu(v).inches if v is not None else 0.0


def px(v):
    return inches(v) * PX


def color_of(fmt):
    try:
        if fmt.type is None:
            return None
        return '#' + str(fmt.fore_color.rgb)
    except Exception:
        return None


def run_color(run):
    try:
        return '#' + str(run.font.color.rgb)
    except Exception:
        return '#333333'


ALIGN = {PP_ALIGN.CENTER: 'center', PP_ALIGN.RIGHT: 'right',
         PP_ALIGN.LEFT: 'left', None: 'left'}
ANCHOR = {MSO_ANCHOR.MIDDLE: 'center', MSO_ANCHOR.BOTTOM: 'flex-end',
          MSO_ANCHOR.TOP: 'flex-start', None: 'flex-start'}


def text_html(tf, default_pt=9.0):
    out = []
    for p in tf.paragraphs:
        runs = []
        for r in p.runs:
            size = r.font.size.pt if r.font.size else default_pt
            style = 'font-size:%.1fpx;line-height:%.1fpx;color:%s;%s' % (
                size * PX / 72.0, size * 1.24 * PX / 72.0, run_color(r),
                'font-weight:700;' if r.font.bold else '')
            runs.append('<span style="%s">%s</span>'
                        % (style, html.escape(r.text)))
        out.append('<div style="text-align:%s">%s</div>'
                   % (ALIGN.get(p.alignment, 'left'), ''.join(runs) or '&nbsp;'))
    return ''.join(out)


def frame_style(tf):
    return 'justify-content:%s' % ANCHOR.get(tf.vertical_anchor, 'flex-start')


def render_table(shape, parts):
    tbl = shape.table
    xs, ys = [px(shape.left)], [px(shape.top)]
    for c in tbl.columns:
        xs.append(xs[-1] + px(c.width))
    for r in tbl.rows:
        ys.append(ys[-1] + px(r.height))

    for ri in range(len(tbl.rows)):
        for ci in range(len(tbl.columns)):
            cell = tbl.cell(ri, ci)
            if cell.is_spanned:
                continue
            w = xs[ci + cell.span_width] - xs[ci]
            h = ys[ri + cell.span_height] - ys[ri]
            fill = color_of(cell.fill)
            style = ('position:absolute;left:%.1fpx;top:%.1fpx;width:%.1fpx;'
                     'height:%.1fpx;box-sizing:border-box;overflow:hidden;'
                     'display:flex;flex-direction:column;%s;'
                     'padding:%.1fpx %.1fpx;%s'
                     % (xs[ci], ys[ri], w, h, frame_style(cell.text_frame),
                        px(cell.margin_top), px(cell.margin_left),
                        'background:%s;' % fill if fill else ''))
            parts.append('<div style="%s">%s</div>'
                         % (style, text_html(cell.text_frame)))
    # grade de referencia, para ver a altura real de cada linha
    for y in ys[1:-1]:
        parts.append('<div style="position:absolute;left:%.1fpx;top:%.1fpx;'
                     'width:%.1fpx;height:1px;background:#E6E6E6"></div>'
                     % (xs[0], y, xs[-1] - xs[0]))


def render_shape(shape, parts):
    if shape.has_table:
        render_table(shape, parts)
        return
    fill = None
    try:
        fill = color_of(shape.fill)
    except Exception:
        pass
    body = ''
    if shape.has_text_frame:
        body = text_html(shape.text_frame)
    style = ('position:absolute;left:%.1fpx;top:%.1fpx;width:%.1fpx;'
             'height:%.1fpx;box-sizing:border-box;display:flex;'
             'flex-direction:column;%s;%s'
             % (px(shape.left), px(shape.top), px(shape.width),
                px(shape.height), frame_style(shape.text_frame)
                if shape.has_text_frame else 'justify-content:flex-start',
                'background:%s;' % fill if fill else ''))
    parts.append('<div style="%s">%s</div>' % (style, body))


def main(src, dst):
    prs = Presentation(src)
    w, h = px(prs.slide_width), px(prs.slide_height)
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            render_shape(shape, parts)
        pages.append('<section style="position:relative;width:%.0fpx;'
                     'height:%.0fpx;background:#fff;margin:0 auto 24px;'
                     'box-shadow:0 2px 12px rgba(0,0,0,.18)">%s'
                     '<div style="position:absolute;right:6px;bottom:2px;'
                     'font:10px Arial;color:#BBB">preview p%d</div>'
                     '</section>' % (w, h, ''.join(parts), i))
    open(dst, 'w', encoding='utf-8').write(
        '<!DOCTYPE html><meta charset="utf-8">'
        '<body style="margin:0;padding:24px;background:#EEE;'
        'font-family:Arial,Helvetica,sans-serif">%s</body>' % ''.join(pages))
    print('preview: %s (%d paginas, %.0fx%.0f px)' % (dst, len(pages), w, h))


if __name__ == '__main__':
    main(sys.argv[1], sys.argv[2])
