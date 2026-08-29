# -*- coding: utf-8 -*-
"""Atualiza o deck ja formatado no template Bain, preservando o branding.

- frente 2: merge com o plano do PI Planning, dividido em duas paginas
  (qualificacao de leads e alocacao de leads), como no proprio cronograma
- frentes 4 e 5: progresso, proximos passos, pontos a escalar e status do mes
- refaz o menu de frentes, que estava quebrando em duas linhas
- marca todo o texto como pt-BR (acentuacao ao editar no PowerPoint)
"""
import shutil, sys, zipfile
from lxml import etree
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from pptx.util import Inches

from deckstyle import (
    STATUS, FONT_NAME, text_w, style_run, WHITE, BLACK, TEXT_DARK, TEXT_REG, TEXT_HELPER, GRAY_LIGHT,
    GRAY_MEDIUM, PLACEHOLDER, CELL_MX, LANG, GREEN_DEEP, NAV_IDLE,
    nlines, lh, split_bold, cell_box, cell_text, set_body_rows, set_row_heights)
import content_update as C

SRC = '/root/.claude/uploads/840b6d74-7f0f-5bc9-b4b2-c44bce458b7b/311204e6-260828__Acompanhamento_das_frentes_v1.pptx'
OUT = '260829__Acompanhamento_das_frentes_v2.pptx'

BODY_STEPS = (10.0, 9.5, 9.0, 8.5, 8.0, 7.5, 7.0, 6.5)
PT_HEAD = 9.0
TABLE_BOTTOM = 7.10          # linha do rodape fica em 7.18


def milestones_table(slide):
    """a maior tabela da pagina e a de milestones"""
    tbls = [sh for sh in slide.shapes if sh.has_table]
    return max(tbls, key=lambda sh: len(sh.table.rows))


def fill_table(shape, rows):
    """reescreve o corpo da tabela de milestones mantendo o estilo do deck"""
    tbl = shape.table
    top = shape.top / 914400
    widths = [c.width / 914400 for c in tbl.columns]
    avail = TABLE_BOTTOM - top

    # colunas que influenciam a altura da linha
    text_cols = [(1, 1), (2, 3), (3, 4), (5, 6), (6, 7), (7, 8)]
    for body_pt in BODY_STEPS:
        id_pt = body_pt - 0.5
        head_h = lh(PT_HEAD) + 0.14
        heights = []
        for r in rows:
            n = nlines(r[0], id_pt, widths[0] - 2 * CELL_MX, True)
            for ci, ri in text_cols:
                n = max(n, nlines(r[ri], body_pt, widths[ci] - 2 * CELL_MX))
            heights.append(max(0.34, n * lh(body_pt) + 0.13))
        if head_h + sum(heights) <= avail:
            break

    # sobra distribuida entre as linhas, com teto para nao inchar
    caps = [max(h, 1.10) for h in heights]
    for _ in range(60):
        slack = avail - (head_h + sum(heights))
        movable = [i for i, h in enumerate(heights) if h < caps[i] - 1e-4]
        if slack <= 0.01 or not movable:
            break
        step = slack / len(movable)
        for i in movable:
            heights[i] = min(caps[i], heights[i] + step)

    set_body_rows(tbl, len(rows))
    set_row_heights(tbl, [head_h] + heights)
    shape.height = int(round((head_h + sum(heights)) * 914400))

    for ri, r in enumerate(rows, start=1):
        mid, desc, bolds, owner, prazo, status, prog, prox, esc = r
        sep = GRAY_MEDIUM if ri == len(rows) else GRAY_LIGHT
        base = (None, None, None, sep)

        c = tbl.cell(ri, 0)
        cell_box(c, None, base, lw=0.5)
        cell_text(c, [(mid, True)], id_pt, TEXT_HELPER)

        for ci, val, col, parts in (
                (1, desc,  TEXT_DARK, split_bold(desc, bolds)),
                (2, owner, TEXT_DARK, [(owner, False)]),
                (3, prazo, TEXT_REG,  [(prazo, False)])):
            c = tbl.cell(ri, ci)
            cell_box(c, None, base, lw=0.5)
            cell_text(c, parts, body_pt, col)

        lbl, bg, fg = STATUS[status]
        c = tbl.cell(ri, 4)
        cell_box(c, bg, (WHITE, WHITE, WHITE, WHITE), lw=3.0)
        cell_text(c, [(lbl, True)], min(body_pt, 8.5), fg,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, mx=0.03)

        for ci, val in ((5, prog), (6, prox), (7, esc)):
            c = tbl.cell(ri, ci)
            cell_box(c, WHITE, (GRAY_MEDIUM,) * 4, lw=0.5)
            colr = PLACEHOLDER if val == C.TBD else TEXT_DARK
            cell_text(c, [(val, False)], body_pt, colr)
    return body_pt


ML, CW = 0.42, 12.4933
NAV_Y, NAV_MAX_PT = 1.52, 10.5


def nav_boxes(slide):
    out = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if len(t) > 3 and t[0].isdigit() and t[1:3] == '. ':
            out.append((int(t[0]), sh, t))
    out.sort(key=lambda x: x[1].left)
    return out


def layout_nav(slide, targets, active):
    """recoloca o menu em uma unica linha, ocupando a largura util"""
    boxes = nav_boxes(slide)
    if len(boxes) != 6:
        return 0
    labels = [t for _, _, t in boxes]
    pt = NAV_MAX_PT
    while pt > 6.5:
        ws = [text_w(t, pt, True) for t in labels]
        gap = (CW - sum(ws)) / (len(labels) - 1)
        if gap >= 0.16:
            break
        pt -= 0.5
    x = ML
    for (n, sh, txt), w in zip(boxes, ws):
        tf = sh.text_frame
        colour = None
        for p in tf.paragraphs:
            for r in p.runs:
                if r.font.color and r.font.color.type is not None:
                    colour = r.font.color.rgb
                    break
            if colour:
                break
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        style_run(p.add_run(), txt, pt,
                  GREEN_DEEP if n == active else (colour or NAV_IDLE), bold=True)
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        sh.left, sh.top = Inches(x), Inches(NAV_Y)
        sh.width, sh.height = Inches(w + 0.06), Inches(lh(pt) + 0.06)
        if n in targets:
            sh.click_action.target_slide = targets[n]
        x += w + gap
    return pt


def set_section_label(slide, text):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip().startswith('Milestones'):
            r = sh.text_frame.paragraphs[0].runs[0]
            r.text = text
            return True
    return False


def delete_slide(prs, index):
    lst = prs.slides._sldIdLst
    el = list(lst)[index]
    prs.part.drop_rel(el.get(qn('r:id')))
    lst.remove(el)


def relink_nav(prs, frente_slides):
    """reaponta o menu de cada pagina para a pagina correta da frente"""
    fixed = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            txt = sh.text_frame.text.strip()
            if len(txt) > 2 and txt[0].isdigit() and txt[1] == '.' and txt[:2] != '1.' or \
               (txt[:3] in ('1. ', '2. ', '3. ', '4. ', '5. ', '6. ')):
                n = int(txt[0])
                if n in frente_slides:
                    sh.click_action.target_slide = frente_slides[n]
                    fixed += 1
    return fixed


def set_lang(path):
    """marca todo o texto como pt-BR, no pacote inteiro"""
    tmp = path + '.tmp'
    shutil.move(path, tmp)
    tags = {qn('a:rPr'), qn('a:endParaRPr'), qn('a:defRPr')}
    touched = 0
    with zipfile.ZipFile(tmp) as zin, \
         zipfile.ZipFile(path, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith('.xml') and (
                    'ppt/slides/' in item.filename or
                    'ppt/slideLayouts/' in item.filename or
                    'ppt/slideMasters/' in item.filename or
                    item.filename == 'ppt/presentation.xml'):
                root = etree.fromstring(data)
                n = 0
                for el in root.iter():
                    if el.tag in tags:
                        el.set('lang', LANG)
                        el.attrib.pop('altLang', None)
                        n += 1
                if n:
                    data = etree.tostring(root, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)
                    touched += n
            zout.writestr(item, data)
    import os; os.remove(tmp)
    return touched


# pagina (indice) -> (frente, milestones, rotulo da secao)
PAGES = [
    (3, 2, C.FRENTE_2A, 'Milestones — Qualificação de leads'),
    (4, 2, C.FRENTE_2B, 'Milestones — Alocação de leads'),
    (7, 4, C.FRENTE_4,  'Milestones'),
    (8, 5, C.FRENTE_5,  'Milestones'),
]
# frente -> pagina de destino do menu
NAV_TARGET = {1: 2, 2: 3, 3: 6, 4: 7, 5: 8, 6: 9}
# todas as paginas de frente, para refazer o menu
FRENTE_PAGES = {2: 1, 3: 2, 4: 2, 6: 3, 7: 4, 8: 5, 9: 6}


def main():
    prs = Presentation(SRC)
    assert 'Coleta de informa' in milestones_table(prs.slides[4]).table.cell(1, 1).text, \
        'a pagina 5 nao e a copia esperada da frente 2'

    pts = {}
    for idx, frente, rows, label in PAGES:
        slide = prs.slides[idx]
        pts[idx] = fill_table(milestones_table(slide), rows)
        set_section_label(slide, label)

    targets = {n: prs.slides[i] for n, i in NAV_TARGET.items()}
    nav_pt = 0
    for idx, active in FRENTE_PAGES.items():
        nav_pt = layout_nav(prs.slides[idx], targets, active) or nav_pt

    prs.save(OUT)
    langs = set_lang(OUT)
    print('tabelas em %s pt | menu em %.1f pt numa linha | %d runs em %s'
          % ('/'.join('%.1f' % pts[i] for i, *_ in PAGES), nav_pt, langs, LANG))


if __name__ == '__main__':
    main()
